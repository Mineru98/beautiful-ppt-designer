#!/usr/bin/env python3
"""
export-ppt.py - Convert PNG slide images into a PowerPoint .pptx file.

Each PNG image becomes a full-bleed slide (image covers the entire slide area).

Requirements:
    pip install python-pptx

Usage:
    python export-ppt.py --input-dir slides/ --output output/presentation.pptx
    python export-ppt.py --input-dir slides/ --output deck.pptx --width 13.333 --height 7.5 --title "My Presentation"
"""

from __future__ import annotations

import argparse
import os
import re
import sys

# ---------------------------------------------------------------------------
# Dependency check
# ---------------------------------------------------------------------------
try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("ERROR: python-pptx is not installed.", file=sys.stderr)
    print("Install it with:  pip install python-pptx", file=sys.stderr)
    sys.exit(1)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _slide_sort_key(filename: str):
    """Sort PNG filenames naturally so slide-2.png comes before slide-10.png."""
    parts = re.split(r"(\d+)", filename)
    return [int(p) if p.isdigit() else p.lower() for p in parts]


def collect_slides(input_dir: str) -> list[str]:
    """Return sorted list of absolute paths to PNG files in *input_dir*."""
    if not os.path.isdir(input_dir):
        raise FileNotFoundError(f"Input directory not found: {input_dir!r}")

    png_files = [f for f in os.listdir(input_dir) if f.lower().endswith(".png")]
    if not png_files:
        raise ValueError(f"No PNG files found in {input_dir!r}")

    png_files.sort(key=_slide_sort_key)
    return [os.path.join(input_dir, f) for f in png_files]


def build_presentation(
    slide_paths: list[str],
    output_path: str,
    width_in: float,
    height_in: float,
    title: str | None,
) -> None:
    """Create a .pptx file where each PNG fills an entire slide."""

    prs = Presentation()

    # Set slide dimensions
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)

    # Blank slide layout (index 6 is universally blank in built-in layouts)
    blank_layout = prs.slide_layouts[6]

    total = len(slide_paths)
    for idx, img_path in enumerate(slide_paths, start=1):
        slide = prs.slides.add_slide(blank_layout)

        # Add image as full-bleed picture (left=0, top=0, fills entire slide)
        slide.shapes.add_picture(
            img_path,
            left=Inches(0),
            top=Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        print(f"  [{idx}/{total}] Added {os.path.basename(img_path)}")

    # Set core properties if a title was provided
    if title:
        prs.core_properties.title = title

    # Ensure output directory exists
    output_dir = os.path.dirname(os.path.abspath(output_path))
    os.makedirs(output_dir, exist_ok=True)

    prs.save(output_path)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def parse_args(argv=None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert PNG slide images exported from Typst into a PowerPoint .pptx file.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument(
        "--input-dir",
        required=True,
        metavar="DIR",
        help="Directory containing PNG slides (e.g. slide-1.png, slide-2.png, ...).",
    )
    parser.add_argument(
        "--output",
        default="output/presentation.pptx",
        metavar="FILE",
        help="Output .pptx file path (default: output/presentation.pptx).",
    )
    parser.add_argument(
        "--width",
        type=float,
        default=13.333,
        metavar="INCHES",
        help="Slide width in inches (default: 13.333 for 16:9).",
    )
    parser.add_argument(
        "--height",
        type=float,
        default=7.5,
        metavar="INCHES",
        help="Slide height in inches (default: 7.5 for 16:9).",
    )
    parser.add_argument(
        "--title",
        default=None,
        metavar="TEXT",
        help="Presentation title stored in document metadata (optional).",
    )
    return parser.parse_args(argv)


def main(argv=None) -> int:
    args = parse_args(argv)

    try:
        print(f"Scanning for PNG slides in: {args.input_dir}")
        slide_paths = collect_slides(args.input_dir)
        print(f"Found {len(slide_paths)} slide(s). Building presentation...")

        build_presentation(
            slide_paths=slide_paths,
            output_path=args.output,
            width_in=args.width,
            height_in=args.height,
            title=args.title,
        )

        print(f"Done. Saved to: {os.path.abspath(args.output)}")
        return 0

    except (FileNotFoundError, ValueError) as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1
    except Exception as exc:  # noqa: BLE001
        print(f"Unexpected error: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
